"""
Jira Bug Analyzer
-----------------
A Streamlit application that ingests a Jira CSV export and produces a
complete bug-analysis dashboard: bug detection, system classification,
repeat-bug clustering (TF-IDF + cosine similarity), account/reporter/
assignee analytics, automated insights, recommendations and a
downloadable text report.

Run:
    pip install -r requirements.txt
    streamlit run app.py
"""

from __future__ import annotations

import io
import re
import zipfile
from datetime import datetime
from typing import Dict, List, Tuple

import numpy as np
import pandas as pd
import streamlit as st
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Page configuration
# ---------------------------------------------------------------------------
st.set_page_config(
    page_title="Jira Bug Analyzer",
    page_icon=":bug:",
    layout="wide",
)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
REQUIRED_COLUMNS = ["Summary", "Issue Type", "Created", "Reporter", "Assignee"]
OPTIONAL_COLUMNS = ["Description", "Account Name", "Seller ID"]

# Keyword maps used for both Bug detection (when Issue Type is missing/ambiguous)
# and for system classification.
BUG_KEYWORDS = [
    "bug", "error", "issue", "fail", "failure", "broken", "crash",
    "not working", "doesn't work", "does not work", "incorrect",
    "wrong", "missing", "unable", "exception", "stuck", "hang",
    "timeout", "blank", "blocker", "defect",
]

SYSTEM_KEYWORDS: Dict[str, List[str]] = {
    # Business / core commerce systems — matched first so they win when
    # keywords overlap with the more general technical buckets below.
    "OMS": [
        "order", "checkout", "payment", "refund", "cart",
        "purchase", "transaction", "order status", "order failed",
    ],
    "IMS": [
        "inventory", "stock", "quantity", "warehouse", "qty",
        "stock mismatch", "oversell",
    ],
    "CMS": [
        "listing", "content", "image", "sku", "catalog",
        "product page", "title", "description missing",
    ],
    "PMS": [
        "discount", "voucher", "promotion", "pricing",
        "promo", "coupon", "price issue",
    ],
    "Chat": [
        "chat", "buyer message", "customer service",
        "buyer reply", "seller reply", "chat not working",
    ],

    # Technical categories — order from most specific to most general so
    # broad terms like "ui" don't swallow more specific matches.
    "Security": [
        "security", "auth", "authentication", "login", "logout",
        "password", "permission", "access denied", "unauthorized",
        "token expired", "sso",
    ],
    "Mobile": [
        "mobile", "android", "ios", "iphone", "ipad",
        "mobile app", "app crash",
    ],
    "Notification": [
        "notification", "alert", "email", "push",
        "sms", "reminder",
    ],
    "Performance": [
        "slow", "latency", "lag", "timeout", "hang",
        "freezing", "loading time", "response time",
    ],
    "Integration": [
        "integration", "api", "webhook", "third party",
        "oauth", "callback",
    ],
    "Dashboard": [
        "dashboard", "report", "analytics", "kpi",
        "graph", "visualization",
    ],
    "Data/Sync": [
        "data mismatch", "not updating", "stale data",
        "data wrong", "sync issue", "sync failed",
    ],
    "UI/UX": [
        "ui", "ux", "alignment", "button", "layout",
        "design", "display", "render", "screen broken",
    ],
    "Backend": [
        "server error", "500 error", "database",
        "backend failure",
    ],
    "Frontend": [
        "frontend", "javascript", "page crash",
    ],
    "Configuration": [
        "config", "setup issue", "misconfiguration",
    ],
    "User Error": [
        "wrong input", "user mistake", "invalid entry",
    ],
}

SIMILARITY_THRESHOLD = 0.55  # cosine similarity to mark two summaries as the same bug


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _safe_lower(value) -> str:
    """Return a lowercase string, treating NaN/None as empty."""
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return ""
    return str(value).lower()


def _contains_any(text: str, keywords: List[str]) -> bool:
    return any(kw in text for kw in keywords)


def validate_columns(df: pd.DataFrame) -> Tuple[bool, List[str]]:
    """Check that required columns exist; return (ok, missing_list)."""
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    return len(missing) == 0, missing


# ---------------------------------------------------------------------------
# Data processing
# ---------------------------------------------------------------------------
def flag_bugs(df: pd.DataFrame) -> pd.DataFrame:
    """Add a Bug_Flag column ('Bug' / 'Non-Bug') using Issue Type + text."""
    df = df.copy()
    issue_type = df["Issue Type"].fillna("").astype(str).str.lower()
    summary = df["Summary"].fillna("").astype(str).str.lower()
    description = (
        df["Description"].fillna("").astype(str).str.lower()
        if "Description" in df.columns
        else pd.Series([""] * len(df), index=df.index)
    )

    combined = (summary + " " + description).str.strip()
    is_bug_type = issue_type.str.contains("bug", na=False)
    is_bug_text = combined.apply(lambda t: _contains_any(t, BUG_KEYWORDS))

    df["Bug_Flag"] = np.where(is_bug_type | is_bug_text, "Bug", "Non-Bug")
    return df


def classify_system(df: pd.DataFrame) -> pd.DataFrame:
    """Add System_Category by matching keywords against summary + description."""
    df = df.copy()
    summary = df["Summary"].fillna("").astype(str).str.lower()
    description = (
        df["Description"].fillna("").astype(str).str.lower()
        if "Description" in df.columns
        else pd.Series([""] * len(df), index=df.index)
    )
    text = summary + " " + description

    def _classify(t: str) -> str:
        for system, keywords in SYSTEM_KEYWORDS.items():
            if _contains_any(t, keywords):
                return system
        return "Others"

    df["System_Category"] = text.apply(_classify)
    return df


def detect_repeats(df: pd.DataFrame, threshold: float = SIMILARITY_THRESHOLD) -> pd.DataFrame:
    """Cluster similar bug summaries with TF-IDF + cosine similarity.

    Adds: Repeat_Count, Cluster_ID. Non-bugs receive Repeat_Count = 1.
    """
    df = df.copy()
    df["Repeat_Count"] = 1
    df["Cluster_ID"] = -1

    bug_mask = df["Bug_Flag"] == "Bug"
    bug_df = df[bug_mask]
    if bug_df.empty:
        return df

    summaries = bug_df["Summary"].fillna("").astype(str).tolist()

    # Edge case: very small corpus
    if len(summaries) < 2:
        df.loc[bug_mask, "Repeat_Count"] = 1
        df.loc[bug_mask, "Cluster_ID"] = 0
        return df

    try:
        vectorizer = TfidfVectorizer(
            stop_words="english",
            ngram_range=(1, 2),
            min_df=1,
            max_features=5000,
        )
        tfidf = vectorizer.fit_transform(summaries)
        sim = cosine_similarity(tfidf)
    except ValueError:
        # All summaries empty / no usable tokens
        df.loc[bug_mask, "Repeat_Count"] = 1
        df.loc[bug_mask, "Cluster_ID"] = 0
        return df

    n = len(summaries)
    cluster_ids = [-1] * n
    next_cluster = 0
    for i in range(n):
        if cluster_ids[i] != -1:
            continue
        cluster_ids[i] = next_cluster
        for j in range(i + 1, n):
            if cluster_ids[j] == -1 and sim[i, j] >= threshold:
                cluster_ids[j] = next_cluster
        next_cluster += 1

    bug_df = bug_df.assign(Cluster_ID=cluster_ids)
    counts = bug_df.groupby("Cluster_ID").size()
    bug_df["Repeat_Count"] = bug_df["Cluster_ID"].map(counts)

    df.loc[bug_mask, "Cluster_ID"] = bug_df["Cluster_ID"].values
    df.loc[bug_mask, "Repeat_Count"] = bug_df["Repeat_Count"].values
    return df


def parse_dates(df: pd.DataFrame) -> pd.DataFrame:
    """Parse the Created column to datetime; add Month for trend charts."""
    df = df.copy()
    df["Created"] = pd.to_datetime(df["Created"], errors="coerce", dayfirst=False)
    df["Month"] = df["Created"].dt.to_period("M").astype(str)
    return df


def get_account_column(df: pd.DataFrame) -> str | None:
    """Return whichever account-identifying column is present, if any."""
    for c in ["Account Name", "Seller ID", "Account", "Customer", "Client"]:
        if c in df.columns:
            return c
    return None


# ---------------------------------------------------------------------------
# Analytics
# ---------------------------------------------------------------------------
def overall_metrics(df: pd.DataFrame) -> Dict[str, float]:
    total = len(df)
    bugs = int((df["Bug_Flag"] == "Bug").sum())
    pct = round((bugs / total * 100), 2) if total else 0.0
    return {"total_tickets": total, "total_bugs": bugs, "bug_pct": pct}


def monthly_trend(df: pd.DataFrame) -> pd.DataFrame:
    bugs = df[df["Bug_Flag"] == "Bug"]
    if bugs.empty or bugs["Month"].isna().all():
        return pd.DataFrame(columns=["Month", "Bugs"])
    trend = (
        bugs.dropna(subset=["Month"])
        .groupby("Month")
        .size()
        .reset_index(name="Bugs")
        .sort_values("Month")
    )
    return trend


def system_breakdown(df: pd.DataFrame) -> pd.DataFrame:
    bugs = df[df["Bug_Flag"] == "Bug"]
    if bugs.empty:
        return pd.DataFrame(columns=["System_Category", "Bugs", "Pct"])
    grouped = (
        bugs.groupby("System_Category")
        .size()
        .reset_index(name="Bugs")
        .sort_values("Bugs", ascending=False)
    )
    grouped["Pct"] = round(grouped["Bugs"] / grouped["Bugs"].sum() * 100, 2)
    return grouped


def top_repeated_bugs(df: pd.DataFrame, n: int = 10) -> pd.DataFrame:
    bugs = df[(df["Bug_Flag"] == "Bug") & (df["Cluster_ID"] >= 0)]
    if bugs.empty:
        return pd.DataFrame(columns=["Sample_Summary", "Repeat_Count", "System_Category"])
    cluster_summary = (
        bugs.sort_values("Repeat_Count", ascending=False)
        .groupby("Cluster_ID")
        .agg(
            Sample_Summary=("Summary", "first"),
            Repeat_Count=("Repeat_Count", "max"),
            System_Category=("System_Category", lambda s: s.mode().iat[0] if not s.mode().empty else "Others"),
        )
        .reset_index(drop=True)
        .sort_values("Repeat_Count", ascending=False)
        .head(n)
    )
    return cluster_summary


def repeated_bugs_over_threshold(df: pd.DataFrame, threshold: int = 3) -> pd.DataFrame:
    repeated = top_repeated_bugs(df, n=10_000)
    return repeated[repeated["Repeat_Count"] > threshold]


def account_breakdown(df: pd.DataFrame, account_col: str | None, n: int = 10) -> pd.DataFrame:
    if not account_col:
        return pd.DataFrame()
    bugs = df[df["Bug_Flag"] == "Bug"]
    if bugs.empty:
        return pd.DataFrame(columns=[account_col, "Bugs"])
    return (
        bugs.groupby(account_col)
        .size()
        .reset_index(name="Bugs")
        .sort_values("Bugs", ascending=False)
        .head(n)
    )


def cross_account_repeats(df: pd.DataFrame, account_col: str | None) -> pd.DataFrame:
    """Find bug clusters that appear across multiple accounts."""
    if not account_col:
        return pd.DataFrame()
    bugs = df[(df["Bug_Flag"] == "Bug") & (df["Cluster_ID"] >= 0)]
    if bugs.empty:
        return pd.DataFrame()
    grouped = (
        bugs.groupby("Cluster_ID")
        .agg(
            Sample_Summary=("Summary", "first"),
            Accounts_Affected=(account_col, lambda s: s.dropna().nunique()),
            Total_Occurrences=("Summary", "size"),
        )
        .reset_index(drop=True)
        .sort_values(["Accounts_Affected", "Total_Occurrences"], ascending=False)
    )
    return grouped[grouped["Accounts_Affected"] > 1].head(15)


def people_breakdown(df: pd.DataFrame, column: str, n: int = 10) -> pd.DataFrame:
    bugs = df[df["Bug_Flag"] == "Bug"]
    if bugs.empty or column not in bugs.columns:
        return pd.DataFrame(columns=[column, "Bugs"])
    return (
        bugs.groupby(column)
        .size()
        .reset_index(name="Bugs")
        .sort_values("Bugs", ascending=False)
        .head(n)
    )


# ---------------------------------------------------------------------------
# Insights & Recommendations
# ---------------------------------------------------------------------------
def generate_insights(
    metrics: Dict[str, float],
    sys_df: pd.DataFrame,
    repeated_df: pd.DataFrame,
    account_df: pd.DataFrame,
    cross_acc_df: pd.DataFrame,
    account_col: str | None,
) -> List[str]:
    insights: List[str] = []

    if metrics["total_tickets"]:
        insights.append(
            f"Out of {metrics['total_tickets']} total tickets, {metrics['total_bugs']} are bugs "
            f"({metrics['bug_pct']}%)."
        )
        if metrics["bug_pct"] >= 40:
            insights.append("Bug ratio is unusually high (>=40%) — quality engineering attention is warranted.")
        elif metrics["bug_pct"] >= 20:
            insights.append("Bug ratio is moderate (20-40%) — track week-over-week to ensure it trends down.")

    if not sys_df.empty:
        top_sys = sys_df.iloc[0]
        insights.append(
            f"Most unstable system: **{top_sys['System_Category']}** with {int(top_sys['Bugs'])} bugs "
            f"({top_sys['Pct']}% of all bugs)."
        )
        if len(sys_df) > 1 and top_sys["Pct"] >= 40:
            insights.append(
                f"{top_sys['System_Category']} alone accounts for >=40% of bugs — consider a dedicated stabilisation sprint."
            )

    if not repeated_df.empty:
        top = repeated_df.iloc[0]
        insights.append(
            f"Top repeated issue appears {int(top['Repeat_Count'])} times — '{top['Sample_Summary'][:120]}'."
        )
        high_repeat = repeated_df[repeated_df["Repeat_Count"] > 3]
        if not high_repeat.empty:
            insights.append(
                f"{len(high_repeat)} bug cluster(s) repeat more than 3 times — likely systemic, not one-off."
            )

    if account_col and not account_df.empty:
        top_acc = account_df.iloc[0]
        insights.append(
            f"Account most affected: **{top_acc[account_col]}** with {int(top_acc['Bugs'])} bug reports."
        )

    if account_col and not cross_acc_df.empty:
        insights.append(
            f"{len(cross_acc_df)} bug cluster(s) span multiple accounts — these are platform-wide, not account-specific."
        )

    if not insights:
        insights.append("No bugs detected in the uploaded dataset.")

    return insights


def generate_recommendations(
    sys_df: pd.DataFrame,
    repeated_df: pd.DataFrame,
    cross_acc_df: pd.DataFrame,
) -> Dict[str, List[str]]:
    short_term: List[str] = []
    long_term: List[str] = []
    operational: List[str] = []

    if not sys_df.empty:
        top_sys = sys_df.iloc[0]["System_Category"]
        short_term.append(
            f"Run a focused triage on open {top_sys} bugs this week and assign a single owner per cluster."
        )
        long_term.append(
            f"Commission an architecture review of {top_sys} — cluster-level repetition suggests a structural issue."
        )

    if not repeated_df.empty:
        high_repeat = repeated_df[repeated_df["Repeat_Count"] > 3]
        if not high_repeat.empty:
            short_term.append(
                "Promote the top repeated bugs into known-issue KB articles with workarounds for support."
            )
            long_term.append(
                "Add automated regression tests for each repeated bug cluster to prevent reintroduction."
            )

    if not cross_acc_df.empty:
        short_term.append("Communicate proactively with affected accounts on the cross-account clusters.")
        long_term.append("Add monitoring/alerting for the top cross-account clusters at the platform layer.")

    operational.extend(
        [
            "Tighten Jira hygiene: enforce Issue Type, Components, and Account fields at ticket creation.",
            "Hold a weekly bug review covering the top 5 repeated clusters and ageing tickets.",
            "Track bug ratio (% of tickets) and repeat-bug ratio as standing engineering KPIs.",
        ]
    )

    if not short_term:
        short_term.append("No immediate fire-fighting items — maintain current quality cadence.")
    if not long_term:
        long_term.append("No structural concerns surfaced — continue routine refactoring.")

    return {"short_term": short_term, "long_term": long_term, "operational": operational}


# ---------------------------------------------------------------------------
# Report builder
# ---------------------------------------------------------------------------
def build_text_report(
    metrics: Dict[str, float],
    sys_df: pd.DataFrame,
    repeated_df: pd.DataFrame,
    account_df: pd.DataFrame,
    insights: List[str],
    recommendations: Dict[str, List[str]],
    account_col: str | None,
) -> str:
    buf = io.StringIO()
    ts = datetime.now().strftime("%Y-%m-%d %H:%M")
    buf.write("JIRA BUG ANALYSIS REPORT\n")
    buf.write(f"Generated: {ts}\n")
    buf.write("=" * 60 + "\n\n")

    buf.write("SUMMARY METRICS\n")
    buf.write("-" * 60 + "\n")
    buf.write(f"Total tickets : {metrics['total_tickets']}\n")
    buf.write(f"Total bugs    : {metrics['total_bugs']}\n")
    buf.write(f"Bug %         : {metrics['bug_pct']}%\n\n")

    buf.write("SYSTEM-WISE BREAKDOWN\n")
    buf.write("-" * 60 + "\n")
    if sys_df.empty:
        buf.write("(no bugs)\n\n")
    else:
        for _, row in sys_df.iterrows():
            buf.write(f"{row['System_Category']:<10}  {int(row['Bugs']):>5} bugs  ({row['Pct']}%)\n")
        buf.write("\n")

    buf.write("TOP REPEATED BUGS\n")
    buf.write("-" * 60 + "\n")
    if repeated_df.empty:
        buf.write("(none)\n\n")
    else:
        for _, row in repeated_df.iterrows():
            summary = str(row["Sample_Summary"])[:120]
            buf.write(f"[{row['System_Category']}] x{int(row['Repeat_Count'])}  {summary}\n")
        buf.write("\n")

    if account_col and not account_df.empty:
        buf.write("TOP AFFECTED ACCOUNTS\n")
        buf.write("-" * 60 + "\n")
        for _, row in account_df.iterrows():
            buf.write(f"{str(row[account_col]):<40}  {int(row['Bugs']):>4} bugs\n")
        buf.write("\n")

    buf.write("INSIGHTS\n")
    buf.write("-" * 60 + "\n")
    for line in insights:
        buf.write(f"- {re.sub(r'[*`]', '', line)}\n")
    buf.write("\n")

    buf.write("RECOMMENDATIONS\n")
    buf.write("-" * 60 + "\n")
    buf.write("Short term:\n")
    for r in recommendations["short_term"]:
        buf.write(f"  - {r}\n")
    buf.write("\nLong term:\n")
    for r in recommendations["long_term"]:
        buf.write(f"  - {r}\n")
    buf.write("\nOperational:\n")
    for r in recommendations["operational"]:
        buf.write(f"  - {r}\n")

    return buf.getvalue()


# ---------------------------------------------------------------------------
# PowerPoint report
# ---------------------------------------------------------------------------
BRAND_BLUE = RGBColor(0x1F, 0x4E, 0x79)
LIGHT_GREY = RGBColor(0x59, 0x59, 0x59)


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Jira Bug Analysis Report"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"


def _add_section_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE


def _add_metrics_slide(prs: Presentation, metrics: Dict[str, float], extra: Dict[str, int]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_section_title(slide, "Summary Metrics")

    rows = [
        ("Total tickets", metrics["total_tickets"]),
        ("Total bugs", metrics["total_bugs"]),
        ("Bug %", f"{metrics['bug_pct']}%"),
        ("Repeated bug clusters", extra.get("repeat_clusters", 0)),
    ]
    left, top = Inches(0.7), Inches(1.6)
    card_w, card_h, gap = Inches(2.9), Inches(2.0), Inches(0.2)
    for i, (label, value) in enumerate(rows):
        x = left + (card_w + gap) * i
        box = slide.shapes.add_textbox(x, top, card_w, card_h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = BRAND_BLUE
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_GREY


def _add_chart_slide(prs: Presentation, title: str, categories: List[str], values: List[float], chart_type) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_title(slide, title)
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    chart_data.add_series("Bugs", values)
    chart = slide.shapes.add_chart(
        chart_type, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7), chart_data
    ).chart
    chart.has_legend = False


def _add_table_slide(prs: Presentation, title: str, df: pd.DataFrame, max_rows: int = 10) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_title(slide, title)
    if df.empty:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
        box.text_frame.text = "No data available."
        return

    df = df.head(max_rows).copy()
    cols = list(df.columns)
    rows, n_cols = len(df) + 1, len(cols)
    tbl_shape = slide.shapes.add_table(rows, n_cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4 * rows + 0.4))
    table = tbl_shape.table
    for j, col in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = BRAND_BLUE
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, col in enumerate(cols):
            value = row[col]
            text = str(value)[:90] if isinstance(value, str) else str(value)
            cell = table.cell(i, j)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)


def _add_bullets_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        clean = re.sub(r"[*`]", "", line)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {clean}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def _add_recommendations_slide(prs: Presentation, recs: Dict[str, List[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_title(slide, "Recommendations")
    headings = [("Short term", recs["short_term"]),
                ("Long term", recs["long_term"]),
                ("Operational", recs["operational"])]
    col_w = Inches(4.1)
    for i, (heading, items) in enumerate(headings):
        x = Inches(0.5) + col_w * i + Inches(0.1 * i)
        box = slide.shapes.add_textbox(x, Inches(1.3), col_w, Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BRAND_BLUE
        for item in items:
            np_ = tf.add_paragraph()
            np_.text = f"•  {item}"
            np_.font.size = Pt(13)
            np_.space_after = Pt(6)


def build_pptx_report(
    metrics: Dict[str, float],
    sys_df: pd.DataFrame,
    repeated_df: pd.DataFrame,
    account_df: pd.DataFrame,
    cross_df: pd.DataFrame,
    rep_df: pd.DataFrame,
    asg_df: pd.DataFrame,
    trend_df: pd.DataFrame,
    insights: List[str],
    recs: Dict[str, List[str]],
    account_col: str | None,
    repeat_clusters: int,
) -> bytes:
    """Return a .pptx report as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    _add_metrics_slide(prs, metrics, {"repeat_clusters": repeat_clusters})

    if not trend_df.empty:
        _add_chart_slide(prs, "Monthly Bug Trend",
                         trend_df["Month"].tolist(),
                         trend_df["Bugs"].tolist(),
                         XL_CHART_TYPE.LINE)

    if not sys_df.empty:
        _add_chart_slide(prs, "Bugs by System",
                         sys_df["System_Category"].tolist(),
                         sys_df["Bugs"].tolist(),
                         XL_CHART_TYPE.BAR_CLUSTERED)
        _add_table_slide(prs, "System Breakdown (table)", sys_df)

    if account_col and not account_df.empty:
        _add_chart_slide(prs, "Top Affected Accounts",
                         account_df[account_col].astype(str).tolist(),
                         account_df["Bugs"].tolist(),
                         XL_CHART_TYPE.BAR_CLUSTERED)

    if not cross_df.empty:
        _add_table_slide(prs, "Bugs Spanning Multiple Accounts", cross_df)

    if not repeated_df.empty:
        _add_table_slide(prs, "Top Repeated Bug Clusters", repeated_df)

    if not rep_df.empty:
        _add_chart_slide(prs, "Top Reporters of Bugs",
                         rep_df["Reporter"].astype(str).tolist(),
                         rep_df["Bugs"].tolist(),
                         XL_CHART_TYPE.BAR_CLUSTERED)
    if not asg_df.empty:
        _add_chart_slide(prs, "Top Assignees of Bugs",
                         asg_df["Assignee"].astype(str).tolist(),
                         asg_df["Bugs"].tolist(),
                         XL_CHART_TYPE.BAR_CLUSTERED)

    _add_bullets_slide(prs, "Auto-generated Insights", insights)
    _add_recommendations_slide(prs, recs)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ---------------------------------------------------------------------------
# Supporting CSV bundle (ZIP)
# ---------------------------------------------------------------------------
def build_supporting_csvs_zip(
    enriched: pd.DataFrame,
    sys_df: pd.DataFrame,
    repeated_df: pd.DataFrame,
    account_df: pd.DataFrame,
    cross_df: pd.DataFrame,
    rep_df: pd.DataFrame,
    asg_df: pd.DataFrame,
    trend_df: pd.DataFrame,
    insights: List[str],
    recs: Dict[str, List[str]],
) -> bytes:
    """Bundle every analysis dataframe + insights/recs into a single ZIP."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("01_enriched_data.csv", enriched.to_csv(index=False))
        zf.writestr("02_system_breakdown.csv", sys_df.to_csv(index=False))
        zf.writestr("03_top_repeated_bugs.csv", repeated_df.to_csv(index=False))
        zf.writestr("04_account_breakdown.csv", account_df.to_csv(index=False) if not account_df.empty else "no account column\n")
        zf.writestr("05_cross_account_clusters.csv", cross_df.to_csv(index=False) if not cross_df.empty else "")
        zf.writestr("06_reporter_breakdown.csv", rep_df.to_csv(index=False))
        zf.writestr("07_assignee_breakdown.csv", asg_df.to_csv(index=False))
        zf.writestr("08_monthly_trend.csv", trend_df.to_csv(index=False))

        ins_csv = pd.DataFrame({"Insight": [re.sub(r"[*`]", "", i) for i in insights]}).to_csv(index=False)
        zf.writestr("09_insights.csv", ins_csv)

        rec_rows = []
        for cat, items in recs.items():
            for it in items:
                rec_rows.append({"Category": cat.replace("_", " ").title(), "Recommendation": it})
        zf.writestr("10_recommendations.csv", pd.DataFrame(rec_rows).to_csv(index=False))

    return buf.getvalue()


# ---------------------------------------------------------------------------
# Caching wrappers (Streamlit re-runs on every interaction)
# ---------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def process(df: pd.DataFrame) -> pd.DataFrame:
    df = parse_dates(df)
    df = flag_bugs(df)
    df = classify_system(df)
    df = detect_repeats(df)
    return df


# ---------------------------------------------------------------------------
# UI
# ---------------------------------------------------------------------------
def main() -> None:
    st.title(":bug: Jira Bug Analyzer")
    st.caption("Upload a Jira CSV export and get an end-to-end bug analysis.")

    with st.expander("Required & optional columns", expanded=False):
        st.markdown(
            "**Required:** `Summary`, `Issue Type`, `Created`, `Reporter`, `Assignee`  \n"
            "**Optional:** `Description`, `Account Name` / `Seller ID`"
        )

    uploaded = st.file_uploader("Upload Jira CSV", type=["csv"])
    if not uploaded:
        st.info("Awaiting CSV upload.")
        return

    # Load
    try:
        raw = pd.read_csv(uploaded)
    except Exception as exc:
        st.error(f"Could not read CSV: {exc}")
        return

    ok, missing = validate_columns(raw)
    if not ok:
        st.error(f"Missing required columns: {', '.join(missing)}")
        st.stop()

    with st.spinner("Crunching tickets…"):
        df = process(raw)

    account_col = get_account_column(df)

    # ---------- Metrics ----------
    metrics = overall_metrics(df)
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Tickets", metrics["total_tickets"])
    c2.metric("Total Bugs", metrics["total_bugs"])
    c3.metric("Bug %", f"{metrics['bug_pct']}%")
    c4.metric("Repeated Clusters", int((df.loc[df["Bug_Flag"] == "Bug", "Repeat_Count"] > 1).sum()))

    st.divider()

    # ---------- Monthly trend ----------
    st.subheader("Monthly Bug Trend")
    trend = monthly_trend(df)
    if trend.empty:
        st.info("No dated bugs available for trend chart.")
    else:
        st.line_chart(trend.set_index("Month"))

    # ---------- System breakdown ----------
    st.subheader("System-wise Bugs")
    sys_df = system_breakdown(df)
    if sys_df.empty:
        st.info("No bugs to break down by system.")
    else:
        col_chart, col_table = st.columns([2, 1])
        col_chart.bar_chart(sys_df.set_index("System_Category")["Bugs"])
        col_table.dataframe(sys_df, use_container_width=True, hide_index=True)

    # ---------- Account breakdown ----------
    st.subheader("Top Affected Accounts")
    if account_col:
        acc_df = account_breakdown(df, account_col)
        if acc_df.empty:
            st.info("No account-level bug data found.")
        else:
            col_chart, col_table = st.columns([2, 1])
            col_chart.bar_chart(acc_df.set_index(account_col)["Bugs"])
            col_table.dataframe(acc_df, use_container_width=True, hide_index=True)

        cross_df = cross_account_repeats(df, account_col)
        if not cross_df.empty:
            st.markdown("**Bug clusters spanning multiple accounts**")
            st.dataframe(cross_df, use_container_width=True, hide_index=True)
    else:
        acc_df = pd.DataFrame()
        cross_df = pd.DataFrame()
        st.caption("No `Account Name` / `Seller ID` column found — skipping account analysis.")

    # ---------- Reporters / Assignees ----------
    st.subheader("Reporters & Assignees")
    col_a, col_b = st.columns(2)
    rep_df = people_breakdown(df, "Reporter")
    asg_df = people_breakdown(df, "Assignee")
    with col_a:
        st.markdown("**Top reporters of bugs**")
        if rep_df.empty:
            st.info("No reporter data.")
        else:
            st.bar_chart(rep_df.set_index("Reporter")["Bugs"])
    with col_b:
        st.markdown("**Top assignees of bugs**")
        if asg_df.empty:
            st.info("No assignee data.")
        else:
            st.bar_chart(asg_df.set_index("Assignee")["Bugs"])

    # ---------- Top repeated bugs ----------
    st.subheader("Top Repeated Bug Clusters")
    repeated_df = top_repeated_bugs(df, n=10)
    if repeated_df.empty:
        st.info("No repeated bug clusters detected.")
    else:
        st.dataframe(repeated_df, use_container_width=True, hide_index=True)
        big = repeated_bugs_over_threshold(df, threshold=3)
        if not big.empty:
            st.warning(f"{len(big)} cluster(s) repeat more than 3 times — likely systemic.")

    # ---------- Insights ----------
    st.subheader("Auto-generated Insights")
    insights = generate_insights(metrics, sys_df, repeated_df, acc_df, cross_df, account_col)
    for line in insights:
        st.markdown(f"- {line}")

    # ---------- Recommendations ----------
    st.subheader("Recommendations")
    recs = generate_recommendations(sys_df, repeated_df, cross_df)
    rc1, rc2, rc3 = st.columns(3)
    with rc1:
        st.markdown("**Short term**")
        for r in recs["short_term"]:
            st.markdown(f"- {r}")
    with rc2:
        st.markdown("**Long term**")
        for r in recs["long_term"]:
            st.markdown(f"- {r}")
    with rc3:
        st.markdown("**Operational**")
        for r in recs["operational"]:
            st.markdown(f"- {r}")

    # ---------- Dataset preview ----------
    with st.expander("Full dataset preview (with derived columns)", expanded=False):
        preview_cols = [
            c for c in [
                "Summary", "Issue Type", "Bug_Flag", "System_Category",
                "Repeat_Count", "Cluster_ID", "Reporter", "Assignee",
                account_col or "", "Created", "Month",
            ] if c and c in df.columns
        ]
        st.dataframe(df[preview_cols], use_container_width=True, hide_index=True)

    # ---------- Export ----------
    st.subheader("Download Reports")

    text_report = build_text_report(
        metrics, sys_df, repeated_df, acc_df, insights, recs, account_col
    )
    enriched_csv = df.to_csv(index=False).encode("utf-8")
    repeat_clusters = int((df.loc[df["Bug_Flag"] == "Bug", "Repeat_Count"] > 1).sum())
    pptx_bytes = build_pptx_report(
        metrics, sys_df, repeated_df, acc_df, cross_df,
        rep_df, asg_df, trend, insights, recs, account_col, repeat_clusters,
    )
    csv_zip_bytes = build_supporting_csvs_zip(
        df, sys_df, repeated_df, acc_df, cross_df,
        rep_df, asg_df, trend, insights, recs,
    )

    d1, d2, d3, d4 = st.columns(4)
    d1.download_button(
        label="PowerPoint deck (.pptx)",
        data=pptx_bytes,
        file_name="jira_bug_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )
    d2.download_button(
        label="Supporting CSVs (.zip)",
        data=csv_zip_bytes,
        file_name="jira_bug_supporting_csvs.zip",
        mime="application/zip",
        use_container_width=True,
    )
    d3.download_button(
        label="Enriched CSV",
        data=enriched_csv,
        file_name="jira_bugs_enriched.csv",
        mime="text/csv",
        use_container_width=True,
    )
    d4.download_button(
        label="Text report",
        data=text_report.encode("utf-8"),
        file_name="jira_bug_report.txt",
        mime="text/plain",
        use_container_width=True,
    )


if __name__ == "__main__":
    main()
